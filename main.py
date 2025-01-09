import pandas as pd
import subprocess
import sys
import comtypes.client
import fitz
import os
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from tqdm import tqdm
import tkinter as tk
from tkinter import filedialog, messagebox

# Função para selecionar arquivo
def select_file(file_type):
    root = tk.Tk()
    root.withdraw()  # Ocultar a janela principal
    file_path = filedialog.askopenfilename(title=f"Selecione o arquivo {file_type}", filetypes=[(file_type, "*.xlsx" if file_type == "Dados.xlsx" else "*.pptx")])
    if not file_path:
        messagebox.showerror("Erro", f"Arquivo {file_type} não selecionado.")
        sys.exit(1)
    return file_path

# Função para selecionar diretório
def select_directory():
    root = tk.Tk()
    root.withdraw()  # Ocultar a janela principal
    directory_path = filedialog.askdirectory(title="Selecione o diretório para salvar o PDF")
    if not directory_path:
        messagebox.showerror("Erro", "Diretório não selecionado.")
        sys.exit(1)
    return os.path.abspath(directory_path)  # Garantir que o caminho seja absoluto

# Selecionar arquivos e diretório
dados_path = select_file("Dados.xlsx")
pptx_template_path = select_file("Certificado_Rondon.pptx")
output_directory = select_directory()

script_dir = os.path.dirname(os.path.abspath(__file__))
os.chdir(script_dir)

# Carregar dados do Excel
df = pd.read_excel(dados_path)
df = df.fillna('')

# Carregar o arquivo PowerPoint
prs_template = Presentation(pptx_template_path)

# Lista para armazenar caminhos dos arquivos PPTX gerados
pptx_files = []

# Gerar certificados em PowerPoint
for index, row in tqdm(df.iterrows(), total=len(df), desc="Gerando certificados em PowerPoint"):
    prs = Presentation(pptx_template_path)  # Carregar o template novamente
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text
                        if any(placeholder in text for placeholder in ['.PESSOA', '.OFICINA1', '.OFICINA2', '.HORAS', '.OPERAÇÃO', '.MUNICIPIO', '.PERIODO', '.COORDENAÇÃO']):
                            text = text.replace('.PESSOA', row['Pessoa'])
                            text = text.replace('.OFICINA1', row['Oficina 1'])
                            text = text.replace('.OFICINA2', row['Oficina 2'])
                            text = text.replace('.HORAS', str(row['Horas']))
                            text = text.replace('.OPERAÇÃO', row['Operação'])
                            text = text.replace('.MUNICIPIO', row['Municipio'])
                            text = text.replace('.PERIODO', row['Periodo'])
                            text = text.replace('.COORDENAÇÃO', row['Coordenação'])
                            run.text = text
    pptx_path = os.path.join(output_directory, f'Certificado_{index}.pptx')
    prs.save(pptx_path)
    pptx_files.append(pptx_path)

# Converter os arquivos PowerPoint para PDF
pdf_files = []

for pptx_file in tqdm(pptx_files, desc="Convertendo PPTX para PDF"):
    pdf_file = pptx_file.replace('.pptx', '.pdf')
    try:
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(pptx_file, WithWindow=False)
        presentation.SaveAs(pdf_file, 32)  # 32 é o formato para PDF
        presentation.Close()
        powerpoint.quit()
        pdf_files.append(pdf_file)
    except Exception as e:
        print(f"Erro ao converter {pptx_file} para PDF: {e}")
        presentation.Close()
        powerpoint.quit()
        sys.exit(1)

# Converter os arquivos PDF para imagens com alta resolução
image_files = []
for pdf_file in tqdm(pdf_files, desc="Convertendo PDF para imagens"):
    doc = fitz.open(pdf_file)
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(dpi=600)  # Ajustar a resolução para 300 DPI
        image_path = pdf_file.replace('.pdf', f'_{page_num}.png')
        pix.save(image_path)
        image_files.append(image_path)
    doc.close()  # Fechar o documento PDF

# Criar o PDF final com dois certificados por página
pdf_path = os.path.join(output_directory, 'Certificados_Final.pdf')
c = canvas.Canvas(pdf_path, pagesize=A4)
width, height = A4

for i, image_file in enumerate(tqdm(image_files, desc="Criando PDF final")):
    if i % 2 == 0 and i != 0:
        c.showPage()
    c.drawImage(image_file, 0, height / 2 if i % 2 == 0 else 0, width, height / 2)

c.save()

# Excluir arquivos temporários
for file_path in pptx_files + pdf_files + image_files:
    try:
        os.remove(file_path)
    except PermissionError:
        print(f"Não foi possível excluir o arquivo: {file_path}")

print(f"\nArquivo PDF salvo em: {pdf_path}\n")